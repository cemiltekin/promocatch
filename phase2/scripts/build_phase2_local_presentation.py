#!/usr/bin/env python3
"""Build the Phase 2 presentation with local screenshots."""

from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parent.parent
DOCS_DIR = PROJECT_ROOT / "docs"
SCREENSHOT_DIR = DOCS_DIR / "screenshots"
OUTPUT_PPTX = DOCS_DIR / "PromoCatch_Project_Explanation_Presentation.pptx"

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(30, 64, 175)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(100, 116, 139)
GREEN = RGBColor(22, 163, 74)
LIGHT_BLUE = RGBColor(239, 246, 255)
LIGHT_GRAY = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
STROKE = RGBColor(203, 213, 225)
FONT = "Aptos"
SLIDE_W = 13.333
SLIDE_H = 7.5


def set_font(paragraph, size: int, color: RGBColor = TEXT, bold: bool = False) -> None:
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.28), Inches(12.1), Inches(0.62))
    frame = title_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = title
    set_font(p, 30, NAVY, True)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.64), Inches(0.93), Inches(11.9), Inches(0.38))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        sp = subtitle_frame.paragraphs[0]
        sp.text = subtitle
        set_font(sp, 13, MUTED)


def add_footer(slide, text: str = "PromoCatch | Software Architecture Project | Phase 2") -> None:
    footer = slide.shapes.add_textbox(Inches(0.62), Inches(7.04), Inches(12.05), Inches(0.25))
    frame = footer.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    set_font(p, 9, MUTED)
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, size: int = 17) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        set_font(p, size, TEXT)
        p.space_after = Pt(10)


def add_panel(slide, left: float, top: float, width: float, height: float, title: str, body: list[str]) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = RGBColor(191, 219, 254)

    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = title
    set_font(p, 18, BLUE, True)
    p.space_after = Pt(8)
    for item in body:
        bp = frame.add_paragraph()
        bp.text = item
        set_font(bp, 13, TEXT)
        bp.space_after = Pt(4)


def add_what_why(slide, what: list[str], why: list[str]) -> None:
    add_panel(slide, 9.05, 1.45, 3.65, 2.45, "What we did", what)
    add_panel(slide, 9.05, 4.15, 3.65, 2.25, "Why it matters", why)


def add_image(slide, image_path: Path, left: float, top: float, width: float, height: float | None = None) -> None:
    if image_path.exists():
        if height is None:
            slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))
        else:
            slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_screenshot_frame(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left - 0.05), Inches(top - 0.05), Inches(width + 0.1), Inches(height + 0.1))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = STROKE
    add_image(slide, image_path, left, top, width, height)


def add_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.95), Inches(1.05), Inches(11.6), Inches(1.0))
    frame = title.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = "PromoCatch"
    set_font(p, 54, WHITE, True)

    subtitle = slide.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(11.3), Inches(1.55))
    sf = subtitle.text_frame
    sf.clear()
    lines = [
        "Campaign and Deal Tracking System",
        "Project Explanation and Phase 2 Demo",
        "Team: Cemil Tekin, Serdar Kaan Kesen, Ömer Tarık Çandır",
        "Local URL used for screenshots: http://127.0.0.1:8010",
    ]
    for idx, line in enumerate(lines):
        sp = sf.paragraphs[0] if idx == 0 else sf.add_paragraph()
        sp.text = line
        set_font(sp, 24 if idx == 0 else 15, RGBColor(226, 232, 240), idx == 0)
        sp.space_after = Pt(8)

    date_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.55), Inches(11.3), Inches(0.3))
    df = date_box.text_frame
    df.clear()
    dp = df.paragraphs[0]
    dp.text = f"Prepared: {datetime.now().strftime('%d.%m.%Y %H:%M')}"
    set_font(dp, 12, RGBColor(203, 213, 225))


def add_project_overview(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Presentation Flow", "The story we will use to explain the project")
    add_bullets(
        slide,
        [
            "1. Explain the problem and the goal of PromoCatch.",
            "2. Introduce users and main campaign management use cases.",
            "3. Show the layered architecture and how each code file fits into it.",
            "4. Summarize what was completed in Phase 2.",
            "5. Walk through local UI, Swagger, and API screenshots.",
            "6. Close with deployment and submission package information.",
        ],
        0.9,
        1.55,
        6.55,
        4.8,
        18,
    )
    add_panel(
        slide,
        7.8,
        1.65,
        4.4,
        3.7,
        "Presentation Goal",
        [
            "Make the architecture understandable",
            "Show that the system runs locally",
            "Connect screenshots to real features",
            "Prepare a clear instructor-facing story",
        ],
    )
    add_footer(slide)


def add_problem_goal(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Problem and Project Goal", "Why PromoCatch exists")
    add_panel(
        slide,
        0.9,
        1.55,
        5.65,
        4.45,
        "Problem",
        [
            "Campaigns are spread across different platforms.",
            "Users need a simple way to view active opportunities.",
            "Campaign data should be searchable, editable, and stored persistently.",
        ],
    )
    add_panel(
        slide,
        6.9,
        1.55,
        5.65,
        4.45,
        "Goal",
        [
            "Build a web-based campaign tracking system.",
            "Provide a frontend for users and a REST API for data operations.",
            "Apply layered architecture and document it using SAD Version 2.",
        ],
    )
    add_footer(slide)


def add_users_use_cases(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Users and Main Use Cases", "What the system allows users to do")
    add_panel(
        slide,
        0.9,
        1.45,
        4.1,
        4.95,
        "Actors",
        [
            "User: tracks, searches, creates, edits, and deletes campaigns.",
            "System: validates data, stores records, and serves API responses.",
        ],
    )
    add_bullets(
        slide,
        [
            "View campaign list",
            "Search campaigns by title or description",
            "Filter campaigns by platform and discount range",
            "Create new campaign records",
            "Update existing campaign records",
            "Delete campaign records",
            "View campaign details through the API",
        ],
        5.45,
        1.55,
        6.9,
        4.9,
        18,
    )
    add_footer(slide)


def add_phase2_summary(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What We Completed in Phase 2", "Small summary of the work added after Phase 1")
    add_bullets(
        slide,
        [
            "Extended SAD Version 2 with completed use case, logical, process, development, and deployment views.",
            "Completed frontend flows: list, search, filter, create, update, and delete campaigns.",
            "Completed backend REST API endpoints for campaign CRUD and health checks.",
            "Added Dockerfile and docker-compose.yml for deployment-ready execution.",
            "Improved Discount Rate input behavior and added fast selection including 100%.",
        ],
        0.9,
        1.5,
        11.7,
        4.9,
        19,
    )
    add_footer(slide)


def add_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Layered Architecture", "How the code maps to the architecture")
    panels = [
        ("Presentation", ["static/index.html", "Forms, campaign cards", "Fetch API client"], 0.75, 1.6),
        ("Control", ["main.py", "routers.py", "REST endpoint definitions"], 3.95, 1.6),
        ("Domain", ["schemas.py", "services.py", "Validation and workflow"], 7.15, 1.6),
        ("Resource", ["database.py", "models.py", "repository.py", "SQLite access"], 10.35, 1.6),
    ]
    for title, body, left, top in panels:
        add_panel(slide, left, top, 2.45, 3.25, title, body)
    add_bullets(
        slide,
        [
            "The browser sends HTTP requests to FastAPI routes.",
            "Routes delegate work to services instead of accessing the database directly.",
            "Repository functions are responsible for SQLAlchemy database operations.",
        ],
        1.0,
        5.2,
        11.2,
        1.1,
        16,
    )
    add_footer(slide)


def add_ui_screenshot(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Local Screen 1: Campaign Management UI", "What we built on the frontend and what it is used for")
    add_screenshot_frame(slide, SCREENSHOT_DIR / "01-home-ui.png", 0.55, 1.35, 8.15, 5.35)
    add_what_why(
        slide,
        [
            "Built a responsive campaign dashboard.",
            "Added search, platform filter, discount range filter, and CRUD actions.",
            "Improved the Discount Rate input and fast selection buttons.",
        ],
        [
            "Users can manage campaigns without using the API manually.",
            "The screen proves the frontend is connected to backend data.",
        ],
    )
    add_footer(slide)


def add_swagger_screenshot(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Local Screen 2: Swagger API Documentation", "What backend endpoints exist and why they are important")
    add_screenshot_frame(slide, SCREENSHOT_DIR / "02-swagger-docs.png", 0.55, 1.35, 8.15, 5.35)
    add_what_why(
        slide,
        [
            "Exposed campaign endpoints with FastAPI.",
            "Documented GET, POST, PUT, DELETE, and health check routes.",
            "Used Pydantic schemas for request and response contracts.",
        ],
        [
            "Swagger makes the backend testable and easy to understand.",
            "It shows the REST API contract used by the frontend.",
        ],
    )
    add_footer(slide)


def add_api_screenshot(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Local Screen 3: Live API Response", "How the backend returns real campaign data")
    add_screenshot_frame(slide, SCREENSHOT_DIR / "03-campaigns-api.png", 0.55, 1.55, 8.15, 4.65)
    add_what_why(
        slide,
        [
            "Returned campaign records as JSON from /campaigns.",
            "Stored and retrieved data through SQLite and SQLAlchemy.",
            "Used the same response format for frontend rendering.",
        ],
        [
            "This proves persistence and backend integration.",
            "It connects the resource layer to the user-facing UI.",
        ],
    )
    add_footer(slide)


def add_deployment(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Deployment and Run Instructions", "Files included for running the Phase 2 project")
    add_bullets(
        slide,
        [
            "Local run: python -m uvicorn main:app --reload",
            "Local test URL used here: http://127.0.0.1:8010",
            "Docker run: docker compose up --build",
            "Persistent data file: data/promocatch.db",
            "Main documentation: docs/SAD_Phase2.pdf and docs/SAD_Phase2.md",
        ],
        0.9,
        1.5,
        11.7,
        4.7,
        19,
    )
    add_footer(slide)


def add_submission_checklist(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Files to Submit to the Instructor", "Recommended Phase 2 delivery package")
    add_bullets(
        slide,
        [
            "1. Phase2_SAD_Code_CemilTekin_SerdarKaanKesen_OmerTarikCandir.zip",
            "2. docs/SAD_Phase2.pdf",
            "3. docs/SAD_Phase2.md",
            "4. docs/PromoCatch_Project_Explanation_Presentation.pptx",
            "5. README.md",
            "6. Source code: main.py, routers.py, services.py, repository.py, models.py, schemas.py, database.py",
            "7. Frontend: static/index.html",
            "8. Deployment files: Dockerfile, docker-compose.yml, requirements.txt",
        ],
        0.9,
        1.45,
        11.8,
        5.2,
        17,
    )
    add_footer(slide)


def add_closing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.fill.background()
    add_title(slide, "Conclusion")
    message = slide.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(10.9), Inches(1.8))
    frame = message.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = "PromoCatch Phase 2 is complete, functional, documented, and ready for presentation."
    set_font(p, 30, GREEN, True)
    p.alignment = PP_ALIGN.CENTER
    add_footer(slide)


def build_presentation() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    add_cover(prs)
    add_project_overview(prs)
    add_problem_goal(prs)
    add_users_use_cases(prs)
    add_phase2_summary(prs)
    add_architecture(prs)
    add_ui_screenshot(prs)
    add_swagger_screenshot(prs)
    add_api_screenshot(prs)
    add_deployment(prs)
    add_submission_checklist(prs)
    add_closing(prs)

    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPTX)
    print(f"Wrote {OUTPUT_PPTX}")


if __name__ == "__main__":
    build_presentation()
